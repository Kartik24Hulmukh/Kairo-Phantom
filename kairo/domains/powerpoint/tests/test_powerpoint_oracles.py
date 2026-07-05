# PROVENANCE: original | PowerPoint domain oracle tests per VERIFICATION_ORACLES.md
"""PowerPoint domain oracle tests — slide_shape_readback + structure_readback + kill-proofs.

Tests verify:
  1. slide_shape_readback: after creating a .pptx, reopen and assert
     shapes/text/tables match the spec. Kill-proof: drop a shape or alter
     text → FAILS.
  2. structure_readback: slide count / layout / table dims survive
     round-trip. Kill-proof: drop a slide → FAILS.
  3. Honest degradation: python-pptx missing → FAIL LOUD.
  4. >=3 gauntlet scenarios: (a) title + bullet-list, (b) table slide,
     (c) image + shape slide — each read-back verified.
  5. Trust stack integration: audit log + egress report.
  6. CLI integration: pptx subcommand works end-to-end.

All tests run fully offline. No mocks on production paths. Zero skips.
"""

from __future__ import annotations

import json
import os
import sys
import tempfile

import pytest
from cryptography.hazmat.primitives.asymmetric import ed25519

_REPO_ROOT = os.path.dirname(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
)
if _REPO_ROOT not in sys.path:
    sys.path.insert(0, _REPO_ROOT)

from kairo.domains.powerpoint.engine import (  # noqa: E402
    create_deck,
    powerpoint_pipeline,
    read_deck,
    save_deck,
)
from kairo.domains.powerpoint.oracles import (  # noqa: E402
    slide_shape_readback,
    structure_readback,
)

# Fixture paths
_FIX = os.path.join(_REPO_ROOT, "kairo", "domains", "powerpoint", "fixtures")
_SPEC_JSON = os.path.join(_FIX, "deck_spec.json")
_TEST_IMAGE = os.path.join(_FIX, "test_image.png")


# ---------------------------------------------------------------------------
# Helper: check engine availability
# ---------------------------------------------------------------------------


def _pptx_available() -> bool:
    try:
        import pptx  # noqa: F401

        return True
    except ImportError:
        return False


_HAS_PPTX = _pptx_available()


# ---------------------------------------------------------------------------
# Helper: load spec and create a deck in a temp file
# ---------------------------------------------------------------------------


def _create_deck_from_spec(spec_path: str, tmpdir: str) -> str:
    """Create a .pptx from a spec file and return the saved path."""
    with open(spec_path, encoding="utf-8") as f:
        spec = json.load(f)

    # Fix image path to be absolute
    for slide in spec.get("slides", []):
        for shape in slide.get("shapes", []):
            if shape.get("type") == "image" and shape.get("path") == "":
                shape["path"] = _TEST_IMAGE

    prs = create_deck(spec)
    output_path = os.path.join(tmpdir, "test_deck.pptx")
    saved_path = save_deck(prs, output_path)
    return saved_path


# ---------------------------------------------------------------------------
# Oracle 1: slide_shape_readback
# ---------------------------------------------------------------------------


class TestSlideShapeReadback:
    """slide_shape_readback oracle: shapes/text/tables match after reopen."""

    def test_text_shapes_match(self):
        """Text shapes read back correctly after save."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available — cannot test slide_shape_readback")

        with tempfile.TemporaryDirectory() as tmp:
            saved_path = _create_deck_from_spec(_SPEC_JSON, tmp)

            # Layouts with placeholders add extra empty shapes.
            # We check only the shapes we created (non-placeholder shapes with text/tables).
            # Slide 0 (Title Slide): 2 placeholders + 2 textboxes = 4 shapes
            # Slide 1 (Title and Content): 2 placeholders + 2 textboxes = 4 shapes
            # Slide 2 (Blank): 1 textbox + 1 table = 2 shapes
            # Slide 3 (Blank): 1 textbox + 1 image + 1 shape = 3 shapes
            deck_info = read_deck(saved_path)

            # Slide 0: verify our 2 textboxes
            textbox_shapes = [s for s in deck_info.slides[0].shapes if "TEXT_BOX" in s.shape_type]
            assert len(textbox_shapes) == 2
            assert textbox_shapes[0].text == "Quarterly Report"
            assert textbox_shapes[1].text == "Q4 2025 Results"

            # Slide 2: verify table
            table_shapes = [s for s in deck_info.slides[2].shapes if s.table_rows > 0]
            assert len(table_shapes) == 1
            assert table_shapes[0].table_rows == 4
            assert table_shapes[0].table_cols == 3

            # Slide 3: verify image + shape
            non_text_shapes = [
                s for s in deck_info.slides[3].shapes if "TEXT_BOX" not in s.shape_type
            ]
            assert len(non_text_shapes) == 2  # image + shape

    def test_table_dims_match(self):
        """Table dimensions read back correctly."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        with tempfile.TemporaryDirectory() as tmp:
            saved_path = _create_deck_from_spec(_SPEC_JSON, tmp)

            # Check table on slide 2
            deck_info = read_deck(saved_path)
            slide2 = deck_info.slides[2]
            table_shape = [s for s in slide2.shapes if s.table_rows > 0][0]
            assert table_shape.table_rows == 4, f"Expected 4 rows, got {table_shape.table_rows}"
            assert table_shape.table_cols == 3, f"Expected 3 cols, got {table_shape.table_cols}"

    def test_kill_proof_wrong_text_fails(self):
        """Kill-proof: alter expected text → oracle FAILS."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        # Use a simple spec with Blank layouts (no placeholder shapes)
        spec = {
            "slides": [
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Correct Title",
                            "left": 1,
                            "top": 1,
                            "width": 8,
                            "height": 2,
                        },
                    ],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as tmp:
            prs = create_deck(spec)
            saved_path = save_deck(prs, os.path.join(tmp, "kill_text.pptx"))

            expected = [[{"text": "WRONG TITLE"}]]
            with pytest.raises(AssertionError, match="text mismatch"):
                slide_shape_readback(saved_path, expected)

    def test_kill_proof_missing_shape_fails(self):
        """Kill-proof: expect fewer shapes than present → oracle FAILS."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        spec = {
            "slides": [
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Title",
                            "left": 1,
                            "top": 1,
                            "width": 8,
                            "height": 2,
                        },
                        {
                            "type": "text",
                            "text": "Body",
                            "left": 1,
                            "top": 3,
                            "width": 8,
                            "height": 2,
                        },
                    ],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as tmp:
            prs = create_deck(spec)
            saved_path = save_deck(prs, os.path.join(tmp, "kill_missing.pptx"))

            # Expect only 1 shape (actually has 2)
            expected = [[{"text": "Title"}]]
            with pytest.raises(AssertionError, match="shape count mismatch"):
                slide_shape_readback(saved_path, expected)

    def test_kill_proof_wrong_table_dims_fails(self):
        """Kill-proof: wrong table dimensions → oracle FAILS."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        spec = {
            "slides": [
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Data",
                            "left": 1,
                            "top": 0.5,
                            "width": 8,
                            "height": 1,
                        },
                        {
                            "type": "table",
                            "rows": 3,
                            "cols": 2,
                            "data": [["A", "B"], ["1", "2"], ["3", "4"]],
                            "left": 1,
                            "top": 2,
                            "width": 6,
                            "height": 3,
                        },
                    ],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as tmp:
            prs = create_deck(spec)
            saved_path = save_deck(prs, os.path.join(tmp, "kill_table.pptx"))

            expected = [[{"text": "Data"}, {"table_rows": 99, "table_cols": 2}]]
            with pytest.raises(AssertionError, match="table_rows mismatch"):
                slide_shape_readback(saved_path, expected)


# ---------------------------------------------------------------------------
# Oracle 2: structure_readback
# ---------------------------------------------------------------------------


class TestStructureReadback:
    """structure_readback oracle: slide count / layout / table dims survive round-trip."""

    def test_slide_count_matches(self):
        """Slide count survives round-trip."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        with tempfile.TemporaryDirectory() as tmp:
            saved_path = _create_deck_from_spec(_SPEC_JSON, tmp)
            passed = structure_readback(saved_path, expected_slide_count=4)
            assert passed

    def test_table_dims_in_structure(self):
        """Table dimensions verified in structure readback."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        with tempfile.TemporaryDirectory() as tmp:
            saved_path = _create_deck_from_spec(_SPEC_JSON, tmp)
            passed = structure_readback(
                saved_path,
                expected_slide_count=4,
                expected_table_dims=[None, None, (4, 3), None],
            )
            assert passed

    def test_kill_proof_wrong_slide_count_fails(self):
        """Kill-proof: wrong slide count → FAILS."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        with tempfile.TemporaryDirectory() as tmp:
            saved_path = _create_deck_from_spec(_SPEC_JSON, tmp)
            with pytest.raises(AssertionError, match="slide count mismatch"):
                structure_readback(saved_path, expected_slide_count=99)

    def test_kill_proof_wrong_table_dims_fails(self):
        """Kill-proof: wrong table dims in structure → FAILS."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        with tempfile.TemporaryDirectory() as tmp:
            saved_path = _create_deck_from_spec(_SPEC_JSON, tmp)
            with pytest.raises(AssertionError, match="table dims"):
                structure_readback(
                    saved_path,
                    expected_slide_count=4,
                    expected_table_dims=[None, None, (99, 3), None],
                )


# ---------------------------------------------------------------------------
# Honest degradation
# ---------------------------------------------------------------------------


class TestHonestDegradation:
    """python-pptx missing → FAIL LOUD, never fake results."""

    def test_pipeline_without_pptx_fails_loud(self):
        """If python-pptx is not installed, pipeline must fail with clear error."""
        if _HAS_PPTX:
            # If python-pptx IS installed, verify the pipeline works
            with open(_SPEC_JSON, encoding="utf-8") as f:
                spec = json.load(f)
            # Fix image path
            for slide in spec.get("slides", []):
                for shape in slide.get("shapes", []):
                    if shape.get("type") == "image" and shape.get("path") == "":
                        shape["path"] = _TEST_IMAGE

            with tempfile.TemporaryDirectory() as tmp:
                result = powerpoint_pipeline(
                    spec=spec,
                    output_path=os.path.join(tmp, "test.pptx"),
                )
                assert result.ok, "Pipeline should succeed when python-pptx is available"
        else:
            result = powerpoint_pipeline(
                spec={"slides": [{"layout": "Blank", "shapes": []}]},
                output_path="/tmp/test.pptx",
            )
            assert not result.ok
            assert "powerpoint engine unavailable" in result.error.lower()


# ---------------------------------------------------------------------------
# Gauntlet scenarios (>=3, zero skips)
# ---------------------------------------------------------------------------


class TestGauntletScenarios:
    ">=3 end-to-end gauntlet scenarios." ""

    def test_scenario_a_title_and_bullets(self):
        """Scenario (a): title + bullet-list deck, read-back verified."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        spec = {
            "slides": [
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Presentation Title",
                            "left": 1,
                            "top": 1,
                            "width": 8,
                            "height": 2,
                            "font_size": 36,
                        },
                    ],
                },
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Agenda",
                            "left": 1,
                            "top": 0.5,
                            "width": 8,
                            "height": 1,
                        },
                        {
                            "type": "bullets",
                            "bullets": ["Introduction", "Main Points", "Conclusion"],
                            "left": 1,
                            "top": 2,
                            "width": 8,
                            "height": 4,
                        },
                    ],
                },
            ]
        }

        with tempfile.TemporaryDirectory() as tmp:
            prs = create_deck(spec)
            saved_path = save_deck(prs, os.path.join(tmp, "scenario_a.pptx"))

            # Read back and verify
            deck_info = read_deck(saved_path)
            assert deck_info.slide_count == 2

            # Slide 0: title (Blank layout = 1 shape, no placeholders)
            assert len(deck_info.slides[0].shapes) == 1
            assert deck_info.slides[0].shapes[0].text == "Presentation Title"

            # Slide 1: agenda + bullets (Blank layout = 2 shapes)
            assert len(deck_info.slides[1].shapes) == 2
            assert deck_info.slides[1].shapes[0].text == "Agenda"
            # Bullets shape: text frame has the bullet text
            bullet_shape = deck_info.slides[1].shapes[1]
            assert "Introduction" in bullet_shape.text
            assert "Main Points" in bullet_shape.text
            assert "Conclusion" in bullet_shape.text

    def test_scenario_b_table_slide(self):
        """Scenario (b): slide with a real table, read-back verified."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        spec = {
            "slides": [
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Data Table",
                            "left": 1,
                            "top": 0.5,
                            "width": 8,
                            "height": 1,
                        },
                        {
                            "type": "table",
                            "rows": 3,
                            "cols": 2,
                            "data": [["A", "B"], ["1", "2"], ["3", "4"]],
                            "left": 1,
                            "top": 2,
                            "width": 6,
                            "height": 3,
                        },
                    ],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as tmp:
            prs = create_deck(spec)
            saved_path = save_deck(prs, os.path.join(tmp, "scenario_b.pptx"))

            # Verify via structure_readback
            passed = structure_readback(
                saved_path,
                expected_slide_count=1,
                expected_table_dims=[(3, 2)],
            )
            assert passed

            # Verify via slide_shape_readback
            expected = [[{"text": "Data Table"}, {"table_rows": 3, "table_cols": 2}]]
            passed = slide_shape_readback(saved_path, expected)
            assert passed

    def test_scenario_c_image_and_shape(self):
        """Scenario (c): slide with an image + shape, read-back verified."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        spec = {
            "slides": [
                {
                    "layout": "Blank",
                    "shapes": [
                        {
                            "type": "text",
                            "text": "Visual Slide",
                            "left": 1,
                            "top": 0.5,
                            "width": 8,
                            "height": 1,
                        },
                        {
                            "type": "image",
                            "path": _TEST_IMAGE,
                            "left": 1,
                            "top": 2,
                            "width": 4,
                            "height": 3,
                        },
                        {
                            "type": "shape",
                            "shape": "ROUNDED_RECTANGLE",
                            "left": 6,
                            "top": 2,
                            "width": 3,
                            "height": 2,
                        },
                    ],
                }
            ]
        }

        with tempfile.TemporaryDirectory() as tmp:
            prs = create_deck(spec)
            saved_path = save_deck(prs, os.path.join(tmp, "scenario_c.pptx"))

            # Read back and verify
            deck_info = read_deck(saved_path)
            assert deck_info.slide_count == 1
            assert len(deck_info.slides[0].shapes) == 3

            # Verify text
            assert deck_info.slides[0].shapes[0].text == "Visual Slide"

            # Verify structure
            passed = structure_readback(saved_path, expected_slide_count=1)
            assert passed


# ---------------------------------------------------------------------------
# Trust stack integration
# ---------------------------------------------------------------------------


class TestTrustStackIntegration:
    """Audit log + zero-egress report integration."""

    def test_pipeline_emits_audit_and_egress(self):
        """Pipeline with private_key emits audit log + egress report."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        with open(_SPEC_JSON, encoding="utf-8") as f:
            spec = json.load(f)
        # Fix image path
        for slide in spec.get("slides", []):
            for shape in slide.get("shapes", []):
                if shape.get("type") == "image" and shape.get("path") == "":
                    shape["path"] = _TEST_IMAGE

        private_key = ed25519.Ed25519PrivateKey.generate()

        with tempfile.TemporaryDirectory() as tmp:
            result = powerpoint_pipeline(
                spec=spec,
                output_path=os.path.join(tmp, "audit_test.pptx"),
                private_key=private_key,
            )
            assert result.ok
            assert result.audit_log_json, "Audit log JSON should be non-empty"
            assert result.egress_report_json, "Egress report JSON should be non-empty"

            from kairo.oracles.ed25519_audit_log import Ed25519AuditLog
            from kairo.oracles.zero_egress_report import (
                report_from_json,
                verify_zero_egress_report,
            )

            public_key = private_key.public_key()
            entries = Ed25519AuditLog.entries_from_json(result.audit_log_json)
            assert len(entries) > 0
            assert Ed25519AuditLog.verify_chain(entries, public_key)

            report = report_from_json(result.egress_report_json)
            assert verify_zero_egress_report(report, public_key)


# ---------------------------------------------------------------------------
# CLI integration
# ---------------------------------------------------------------------------


class TestCLIIntegration:
    """pptx CLI subcommand works end-to-end via registry."""

    def test_cli_create(self):
        """`kairo pptx create` produces output + audit artifacts."""
        if not _HAS_PPTX:
            pytest.fail("python-pptx not available")

        from kairo.cli import main

        # Create a spec with absolute image path
        with tempfile.TemporaryDirectory() as tmp:
            spec = {
                "slides": [
                    {
                        "layout": "Blank",
                        "shapes": [
                            {
                                "type": "text",
                                "text": "CLI Test Slide",
                                "left": 1,
                                "top": 1,
                                "width": 8,
                                "height": 2,
                            }
                        ],
                    }
                ]
            }
            spec_path = os.path.join(tmp, "spec.json")
            with open(spec_path, "w", encoding="utf-8") as f:
                json.dump(spec, f)

            out_dir = os.path.join(tmp, "pptx_output")
            rc = main(["pptx", "create", spec_path, "--out", "test_cli.pptx", "--outdir", out_dir])
            assert rc == 0, f"CLI create failed with exit code {rc}"
            assert os.path.isfile(os.path.join(out_dir, "audit_log.json"))
            assert os.path.isfile(os.path.join(out_dir, "zero_egress_report.json"))
            assert os.path.isfile(os.path.join(out_dir, "test_cli.pptx"))
